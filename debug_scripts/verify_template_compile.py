"""Compile a varied mock deck against the VI EC template to verify layout fidelity."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

from presentations.agents.assembler import compile_from_template
from presentations.core.schemas import DeckSpec, GenerationMode, PlaceholderMapping, SlideSpec
from presentations.ingest.pptx_layout import generate_layout_map

TEMPLATE = Path("documentation/briefs/VI_EC_Corporate_PPT_Template_2026.pptx")
OUTPUT = Path("debug_scripts/template_compile_verify.pptx")


def _first_layout_with_role(profile, role: str) -> int:
    for idx, entry in sorted(profile.layouts.items()):
        if entry.role == role and entry.placeholders:
            return idx
    raise ValueError(f"No layout with role {role}")


def main() -> int:
    if not TEMPLATE.exists():
        print(f"Template missing: {TEMPLATE}", file=sys.stderr)
        return 1

    profile = generate_layout_map(TEMPLATE)
    title_idx = _first_layout_with_role(profile, "title")
    section_idx = _first_layout_with_role(profile, "section")
    two_idx = _first_layout_with_role(profile, "two-content")
    picture_idx = _first_layout_with_role(profile, "picture")

    def ph(layout_idx: int, prefer_title: bool = False) -> int:
        entry = profile.layouts[layout_idx]
        if prefer_title:
            for ph_info in entry.placeholders:
                if "TITLE" in ph_info.type.upper():
                    return ph_info.index
        for ph_info in entry.placeholders:
            if "BODY" in ph_info.type.upper() or "TITLE" in ph_info.type.upper():
                return ph_info.index
        return entry.placeholders[0].index

    deck = DeckSpec(
        title="MiFID II Investor Briefing",
        mode=GenerationMode.TEMPLATE,
        slides=[
            SlideSpec(
                layout_index=title_idx,
                mappings=[PlaceholderMapping(ph_idx=ph(title_idx, True), content="Impact of MiFID II on EU Retail Investors")],
            ),
            SlideSpec(
                layout_index=section_idx,
                mappings=[PlaceholderMapping(ph_idx=ph(section_idx, True), content="Key Requirements")],
            ),
            SlideSpec(
                layout_index=two_idx,
                mappings=[
                    PlaceholderMapping(ph_idx=ph(two_idx, True), content="Transparency vs Best Execution"),
                    PlaceholderMapping(ph_idx=ph(two_idx), content="Clear disclosure rules\nBest execution obligations"),
                ],
            ),
            SlideSpec(
                layout_index=picture_idx,
                mappings=[PlaceholderMapping(ph_idx=ph(picture_idx, True), content="Investor Protection Measures")],
            ),
        ],
    )

    compile_from_template(TEMPLATE, deck, OUTPUT, profile)
    prs = Presentation(str(OUTPUT))
    layout_indices = {slide.slide_layout.name for slide in prs.slides}
    print(f"Compiled {len(prs.slides)} slides to {OUTPUT}")
    print(f"Distinct layout names: {len(layout_indices)}")
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if "**" in text:
                print("FAIL: markdown leaked in output", text[:80])
                return 1
    print("OK: no markdown leakage detected")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
