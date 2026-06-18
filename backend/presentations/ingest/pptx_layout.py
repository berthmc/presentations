"""Extract layout profiles from .pptx templates."""

from pathlib import Path

from pptx import Presentation

from presentations.compile.theme_palette import extract_theme_metadata
from presentations.core.layout_roles import classify_layout_role, detect_media_flags, summarize_layout
from presentations.core.schemas import LayoutEntry, LayoutProfile, PlaceholderInfo


def _read_geometry(shape) -> tuple[int | None, int | None, int | None, int | None]:
    """Read placeholder geometry in EMU, returning None for inherited values."""
    try:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
    except AttributeError:
        return None, None, None, None
    return left, top, width, height


def _build_placeholder_info(shape) -> PlaceholderInfo:
    """Build PlaceholderInfo with geometry from a layout placeholder shape."""
    left, top, width, height = _read_geometry(shape)
    return PlaceholderInfo(
        index=shape.placeholder_format.idx,
        name=shape.name,
        type=str(shape.placeholder_format.type),
        left=left,
        top=top,
        width=width,
        height=height,
    )


def generate_layout_map(template_path: str | Path) -> LayoutProfile:
    """Analyze master slides in a target presentation template and return a mapping profile.

    Args:
        template_path: Path to the corporate .pptx template.

    Returns:
        LayoutProfile with layout indices, placeholder metadata, roles, and theme.
    """
    path = Path(template_path)
    prs = Presentation(str(path))
    layout_profile: dict[int, LayoutEntry] = {}

    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = [_build_placeholder_info(ph) for ph in layout.placeholders]
        draft = LayoutEntry(name=layout.name, placeholders=placeholders)
        has_picture, has_chart, has_table = detect_media_flags(draft)
        role = classify_layout_role(draft)
        summary = summarize_layout(draft)
        layout_profile[idx] = LayoutEntry(
            name=layout.name,
            placeholders=placeholders,
            role=role,
            summary=summary,
            has_picture=has_picture,
            has_chart=has_chart,
            has_table=has_table,
        )

    theme = extract_theme_metadata(prs)
    return LayoutProfile(
        source_path=str(path.resolve()),
        source_type="pptx",
        layouts=layout_profile,
        theme=theme,
    )
