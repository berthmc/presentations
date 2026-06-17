"""Stage 5: Constraint validation loop (Inspector)."""

import asyncio
from pathlib import Path

from loguru import logger
from markitdown import MarkItDown
from pptx import Presentation

from presentations.agents.skill_rules import (
    VISUAL_QA_PROMPT,
    build_rollback_feedback,
    check_leftover_placeholder_text,
    check_unicode_bullets,
    estimate_text_overflow,
)
from presentations.config.settings import get_settings
from presentations.core.schemas import QAIssue, QAReport
from presentations.core.state import PipelineState
from presentations.llm.router import LLMRouter
from presentations.qa.geometric import check_geometric
from presentations.qa.render import render_slides_to_images


async def _audit_content_text(pptx_path: Path) -> list[QAIssue]:
    """Run markitdown content QA per SKILL.md."""
    issues: list[QAIssue] = []
    try:
        md = MarkItDown()
        result = await asyncio.to_thread(md.convert, str(pptx_path))
        text = result.text_content or ""
    except Exception as exc:
        logger.warning("Content QA markitdown failed: {}", exc)
        return issues

    for message in check_leftover_placeholder_text(text):
        issues.append(QAIssue(slide=0, severity="error", category="content", message=message))
    return issues


def _audit_placeholder_overflow(pptx_path: Path, deck_spec) -> list[QAIssue]:
    """Estimate text overflow in filled placeholders."""
    issues: list[QAIssue] = []
    prs = Presentation(str(pptx_path))
    for slide_num, slide_spec in enumerate(deck_spec.slides, start=1):
        if slide_num > len(prs.slides):
            break
        slide = prs.slides[slide_num - 1]
        for mapping in slide_spec.mappings:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx != mapping.ph_idx:
                    continue
                width = shape.width
                height = shape.height
                estimate = estimate_text_overflow(
                    mapping.content,
                    width_emu=width,
                    height_emu=height,
                )
                if estimate["overflow"]:
                    issues.append(
                        QAIssue(
                            slide=slide_num,
                            severity="error",
                            category="overflow",
                            message=(
                                f"Placeholder ph_idx={mapping.ph_idx} may overflow by "
                                f"{estimate['overflow_lines']} line(s) "
                                f"({estimate['total_lines']}/{estimate['max_lines']} lines)"
                            ),
                        )
                    )
                for bullet_issue in check_unicode_bullets(mapping.content):
                    issues.append(
                        QAIssue(
                            slide=slide_num,
                            severity="warning",
                            category="formatting",
                            message=bullet_issue,
                        )
                    )
    return issues


async def _audit_visual(
    image_paths: list[Path],
    *,
    allow_cloud: bool,
    vlm_enabled: bool,
) -> tuple[list[QAIssue], list[str]]:
    """Run geometric + optional VLM visual QA using SKILL.md prompt."""
    issues: list[QAIssue] = []
    reasons: list[str] = []

    for index, image_path in enumerate(image_paths, start=1):
        issues.extend(check_geometric(image_path, index))

    if not vlm_enabled:
        return issues, reasons

    router = LLMRouter(allow_cloud=allow_cloud)
    vision = await router.get_vision_provider()
    if vision is None:
        return issues, reasons

    for index, image_path in enumerate(image_paths, start=1):
        try:
            result = await vision.audit_slide_image(str(image_path), VISUAL_QA_PROMPT)
            if not result.get("passed", False):
                for reason in result.get("reasons", []):
                    reasons.append(str(reason))
                    issues.append(
                        QAIssue(slide=index, severity="error", category="vlm", message=str(reason))
                    )
        except Exception as exc:
            logger.warning("VLM audit failed for slide {}: {}", index, exc)
            issues.append(
                QAIssue(
                    slide=index,
                    severity="warning",
                    category="vlm",
                    message=f"VLM audit unavailable: {exc}",
                )
            )
    return issues, reasons


async def run_inspector(state: PipelineState) -> PipelineState:
    """Validate compiled deck; populate qa_report and rollback_reasons on failure.

    Args:
        state: Pipeline state with output_path and deck_spec.

    Returns:
        Updated state with qa_report; rollback_reasons set when validation fails.
    """
    settings = get_settings()
    if not state.request.run_qa or not state.output_path:
        state.qa_report = QAReport(passed=True, reasons=["QA skipped"])
        return state

    pptx_path = Path(state.output_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"Compiled deck not found: {pptx_path}")

    all_issues: list[QAIssue] = []
    all_issues.extend(await _audit_content_text(pptx_path))
    if state.deck_spec:
        all_issues.extend(_audit_placeholder_overflow(pptx_path, state.deck_spec))

    slide_images: list[str] = []
    try:
        images = await render_slides_to_images(pptx_path)
        slide_images = [str(p.resolve()) for p in images]
        vlm_enabled = settings.qa_vlm_enabled or state.request.allow_cloud
        visual_issues, _visual_reasons = await _audit_visual(
            images,
            allow_cloud=state.request.allow_cloud,
            vlm_enabled=vlm_enabled,
        )
        all_issues.extend(visual_issues)
    except FileNotFoundError as exc:
        logger.warning("Visual QA skipped: {}", exc)
        all_issues.append(
            QAIssue(slide=0, severity="warning", category="render", message=f"QA render skipped: {exc}")
        )

    passed = not any(issue.severity == "error" for issue in all_issues)
    reasons = [issue.message for issue in all_issues if issue.severity == "error"]
    state.qa_report = QAReport(
        passed=passed,
        reasons=reasons,
        issues=all_issues,
        slide_images=slide_images,
        iterations=state.revision + 1,
    )
    if not passed:
        state.rollback_reasons = build_rollback_feedback(reasons)
    return state
