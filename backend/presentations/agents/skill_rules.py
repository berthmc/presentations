"""Formatting and compliance rules distilled from presentation skill docs.

Sources:
- documentation/briefs/editing.md
- documentation/briefs/pptxgenjs.md
- documentation/briefs/SKILL.md
"""

import re
from dataclasses import dataclass
from typing import Any

from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN

# --- Content patterns (SKILL.md + editing.md) ---

LEFTOVER_PLACEHOLDER_PATTERN = re.compile(
    r"xxxx|lorem|ipsum|this.*(page|slide).*layout",
    re.IGNORECASE,
)
UNICODE_BULLET_PATTERN = re.compile(r"^[\u2022\u2023\u2043\u2219\u25CF\u25E6\u2024\-]\s*")
LEADING_MARKDOWN_PATTERN = re.compile(r"^(#+\s+|-\s+|\*\s+|\d+\.\s+)")
INLINE_MARKDOWN_PATTERN = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*")
LEADIN_LABEL_PATTERN = re.compile(r"^([^:]{1,40}):\s*(.+)$")
SMART_QUOTE_MAP = {
    "\u201c": "&#x201C;",
    "\u201d": "&#x201D;",
    "\u2018": "&#x2018;",
    "\u2019": "&#x2019;",
}

# --- Visual QA prompt (SKILL.md) ---

VISUAL_QA_PROMPT = """Visually inspect these slides. Assume there are issues — find them.

Look for:
- Overlapping elements (text through shapes, lines through words, stacked elements)
- Text overflow or cut off at edges/box boundaries
- Decorative lines positioned for single-line text but title wrapped to two lines
- Source citations or footers colliding with content above
- Elements too close (< 0.3" gaps) or cards/sections nearly touching
- Uneven gaps (large empty area in one place, cramped in another)
- Insufficient margin from slide edges (< 0.5")
- Columns or similar elements not aligned consistently
- Low-contrast text (e.g., light gray text on cream-colored background)
- Low-contrast icons (e.g., dark icons on dark backgrounds without a contrasting circle)
- Text boxes too narrow causing excessive wrapping
- Leftover placeholder content

For each slide, list issues or areas of concern, even if minor.

Return strict JSON: {"passed": boolean, "reasons": ["..."]}
"""


@dataclass(frozen=True)
class TextSegment:
    """One inline text segment with optional emphasis."""

    text: str
    bold: bool = False
    italic: bool = False
    accent: bool = False


def strip_leading_markdown(line: str) -> str:
    """Remove leading markdown bullets, numbering, or heading markers."""
    cleaned = UNICODE_BULLET_PATTERN.sub("", line.strip())
    return LEADING_MARKDOWN_PATTERN.sub("", cleaned).strip()


def normalize_content_lines(content: str) -> list[str]:
    """Split content into lines, stripping unicode bullets per editing.md."""
    lines: list[str] = []
    for raw in content.split("\n"):
        line = strip_leading_markdown(raw)
        if line:
            lines.append(line)
    return lines


def parse_inline_markdown(line: str) -> list[TextSegment]:
    """Parse inline **bold** and *italic* markers into text segments."""
    if "**" not in line and line.count("*") < 2:
        lead_in = _parse_leadin_label(line)
        if lead_in is not None:
            return lead_in

    segments: list[TextSegment] = []
    cursor = 0
    for match in INLINE_MARKDOWN_PATTERN.finditer(line):
        if match.start() > cursor:
            segments.append(TextSegment(text=line[cursor : match.start()]))

        if match.group(1) is not None:
            segments.append(TextSegment(text=match.group(1), bold=True, accent=True))
        elif match.group(2) is not None:
            segments.append(TextSegment(text=match.group(2), italic=True))

        cursor = match.end()

    if cursor < len(line):
        segments.append(TextSegment(text=line[cursor:]))

    if not segments:
        return [TextSegment(text=line)]
    return [segment for segment in segments if segment.text]


def _parse_leadin_label(line: str) -> list[TextSegment] | None:
    """Split ``Label: body`` lines into bold label + normal body segments."""
    match = LEADIN_LABEL_PATTERN.match(line)
    if match is None:
        return None
    return [
        TextSegment(text=f"{match.group(1)}:", bold=True, accent=True),
        TextSegment(text=match.group(2)),
    ]


def _apply_accent_to_run(run, accent: MSO_THEME_COLOR | None) -> None:
    """Apply a theme accent colour when supported by python-pptx."""
    if accent is None:
        return
    try:
        run.font.color.theme_color = accent
    except AttributeError:
        return


def _apply_font_to_run(run, font_name: str | None) -> None:
    """Apply a template font name to a text run when provided."""
    if not font_name:
        return
    try:
        run.font.name = font_name
    except AttributeError:
        return


def _add_segments_to_paragraph(
    paragraph,
    segments: list[TextSegment],
    *,
    is_title: bool,
    accent: MSO_THEME_COLOR | None,
    tint_title_accent: bool,
    font_name: str | None = None,
) -> None:
    """Add parsed inline segments to a paragraph as separate runs."""
    for segment in segments:
        run = paragraph.add_run()
        run.text = segment.text
        run.font.bold = is_title or segment.bold
        run.font.italic = segment.italic
        _apply_font_to_run(run, font_name)

        if is_title and tint_title_accent and accent is not None:
            _apply_accent_to_run(run, accent)
        elif segment.accent and accent is not None:
            _apply_accent_to_run(run, accent)


def fill_placeholder_with_rules(
    placeholder,
    content: str,
    *,
    is_title: bool = False,
    accent: MSO_THEME_COLOR | None = None,
    tint_title_accent: bool = False,
    font_name: str | None = None,
) -> None:
    """Fill a python-pptx placeholder applying editing.md formatting rules."""
    lines = normalize_content_lines(content)
    if not lines:
        placeholder.text = ""
        return

    text_frame = placeholder.text_frame
    text_frame.clear()

    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT

        if is_title and "**" not in line and "*" not in line:
            run = paragraph.add_run()
            run.text = line
            run.font.bold = True
            _apply_font_to_run(run, font_name)
            if tint_title_accent and accent is not None:
                _apply_accent_to_run(run, accent)
            continue

        segments = parse_inline_markdown(line)
        if len(segments) == 1 and not segments[0].bold and not segments[0].italic:
            run = paragraph.add_run()
            run.text = segments[0].text
            run.font.bold = is_title or (index == 0 and len(lines) > 1 and ":" in line)
            _apply_font_to_run(run, font_name)
            if run.font.bold and accent is not None and not is_title:
                _apply_accent_to_run(run, accent)
            elif is_title and tint_title_accent and accent is not None:
                _apply_accent_to_run(run, accent)
            continue

        _add_segments_to_paragraph(
            paragraph,
            segments,
            is_title=is_title,
            accent=accent,
            tint_title_accent=tint_title_accent,
            font_name=font_name,
        )


def estimate_text_overflow(
    text: str,
    *,
    width_emu: int,
    height_emu: int,
    font_size_pt: float = 14.0,
) -> dict[str, Any]:
    """Estimate whether text exceeds placeholder bounds (Inspector heuristic)."""
    emu_per_inch = 914400
    width_in = width_emu / emu_per_inch
    height_in = height_emu / emu_per_inch
    chars_per_line = max(1, int(width_in * (72 / font_size_pt) * 0.55))
    line_height_in = font_size_pt / 72 * 1.3
    max_lines = max(1, int(height_in / line_height_in))

    lines = text.split("\n") if text else [""]
    total_lines = 0
    for line in lines:
        total_lines += max(1, (len(line) + chars_per_line - 1) // chars_per_line)

    overflow_lines = max(0, total_lines - max_lines)
    return {
        "total_lines": total_lines,
        "max_lines": max_lines,
        "overflow_lines": overflow_lines,
        "overflow": overflow_lines > 0,
        "chars_per_line": chars_per_line,
    }


def check_leftover_placeholder_text(text: str) -> list[str]:
    """Return issues for leftover template placeholder text (SKILL.md content QA)."""
    issues: list[str] = []
    for line in text.splitlines():
        if LEFTOVER_PLACEHOLDER_PATTERN.search(line):
            issues.append(f"Leftover placeholder text detected: {line[:80]}")
    return issues


def check_unicode_bullets(content: str) -> list[str]:
    """Return issues when unicode bullets are used instead of layout bullets."""
    issues: list[str] = []
    for line in content.splitlines():
        if UNICODE_BULLET_PATTERN.match(line.strip()):
            issues.append("Unicode bullet character detected; use layout bullets instead")
    return issues


def build_rollback_feedback(issues: list[str]) -> list[str]:
    """Format Inspector findings for Planner rollback."""
    return [issue for issue in issues if issue.strip()]
