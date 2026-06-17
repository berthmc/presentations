"""Stage 4: Programmatic compilation into native .pptx (Assembler)."""

import time
from pathlib import Path
from uuid import uuid4

from loguru import logger
from pptx import Presentation

from presentations.agents.skill_rules import (
    check_leftover_placeholder_text,
    fill_placeholder_with_rules,
    normalize_content_lines,
)
from presentations.compile.pptxgen_runner import compile_from_scratch
from presentations.config.pipeline_logging import summarize_deck_spec, summarize_pptx_file
from presentations.core.schemas import GenerationMode
from presentations.core.state import PipelineState
from presentations.llm.layout_validate import sanitize_deck_spec


def _find_placeholder(slide, ph_idx: int):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            return shape
    return None


def _is_title_placeholder(name: str) -> bool:
    lowered = name.lower()
    return "title" in lowered and "subtitle" not in lowered


def compile_from_template(
    template_path: str | Path,
    deck_spec,
    output_path: str | Path,
    layout_profile=None,
) -> Path:
    """Compile a deck by filling template placeholders with skill-rule formatting."""
    template = Path(template_path)
    destination = Path(output_path)
    destination.parent.mkdir(parents=True, exist_ok=True)

    if layout_profile is not None:
        original_slides = len(deck_spec.slides)
        deck_spec = sanitize_deck_spec(deck_spec, layout_profile)
        if len(deck_spec.slides) != original_slides:
            logger.info("Sanitized deck spec: {} slides after layout clamping", len(deck_spec.slides))

    prs = Presentation(str(template))
    slide_ids = list(prs.slides._sldIdLst)  # noqa: SLF001
    for slide_id in slide_ids:
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001

    for slide_num, slide_spec in enumerate(deck_spec.slides, start=1):
        layout = prs.slide_layouts[slide_spec.layout_index]
        slide = prs.slides.add_slide(layout)
        filled_ph: set[int] = set()
        content_chars = sum(len(mapping.content) for mapping in slide_spec.mappings)
        logger.debug(
            "Slide {} layout_index={} mappings={} content_chars={}",
            slide_num,
            slide_spec.layout_index,
            len(slide_spec.mappings),
            content_chars,
        )
        for mapping in slide_spec.mappings:
            placeholder = _find_placeholder(slide, mapping.ph_idx)
            if placeholder is None:
                logger.warning(
                    "Placeholder {} not found on layout {}",
                    mapping.ph_idx,
                    slide_spec.layout_index,
                )
                continue
            ph_name = placeholder.name or ""
            fill_placeholder_with_rules(
                placeholder,
                mapping.content,
                is_title=_is_title_placeholder(ph_name),
            )
            filled_ph.add(mapping.ph_idx)

        if layout_profile and slide_spec.layout_index in layout_profile.layouts:
            entry = layout_profile.layouts[slide_spec.layout_index]
            for ph in entry.placeholders:
                if ph.index in filled_ph:
                    continue
                shape = _find_placeholder(slide, ph.index)
                if shape is None:
                    continue
                leftover = check_leftover_placeholder_text(shape.text or "")
                if leftover or not normalize_content_lines(shape.text):
                    sp = shape._element  # noqa: SLF001
                    sp.getparent().remove(sp)

    prs.save(str(destination))
    return destination


async def run_assembler(state: PipelineState) -> PipelineState:
    """Compile deck_spec into a native .pptx binary.

    Args:
        state: Pipeline state with deck_spec and template_path.

    Returns:
        Updated state with output_path.
    """
    if state.deck_spec is None:
        raise ValueError("Assembler requires deck_spec")

    from presentations.config.settings import get_settings

    settings = get_settings()
    deck_summary = summarize_deck_spec(state.deck_spec)
    output_name = f"{state.deck_spec.title.replace(' ', '_')}_{uuid4().hex[:8]}.pptx"
    output_path = settings.output_dir / output_name
    mode = GenerationMode(state.mode or state.request.mode.value)

    if mode == GenerationMode.TEMPLATE:
        logger.info(
            "Compiling template deck slides={} mappings={} template={}",
            deck_summary["slides"],
            deck_summary["mappings"],
            state.template_path,
        )
        if not state.template_path:
            raise ValueError("template_path required for template mode")
        compile_started = time.perf_counter()
        compile_from_template(
            state.template_path,
            state.deck_spec,
            output_path,
            state.layout_profile,
        )
        compile_duration = time.perf_counter() - compile_started
        compile_mode = "template"
    else:
        theme_name = "tech"
        if state.layout_profile and state.layout_profile.theme:
            theme_name = state.layout_profile.theme.get("brand", "tech")
        logger.info(
            "Compiling scratch deck slides={} mappings={} theme={}",
            deck_summary["slides"],
            deck_summary["mappings"],
            theme_name,
        )
        compile_started = time.perf_counter()
        await compile_from_scratch(state.deck_spec, output_path, theme_name=theme_name)
        compile_duration = time.perf_counter() - compile_started
        compile_mode = "scratch"

    file_summary = summarize_pptx_file(output_path)
    logger.info(
        "Assembler compiled {} deck slides={} mappings={} size_kb={} path={} duration_s={:.2f}",
        compile_mode,
        deck_summary["slides"],
        deck_summary["mappings"],
        file_summary["size_kb"],
        file_summary["path"],
        compile_duration,
    )

    state.output_path = str(output_path.resolve())
    return state
