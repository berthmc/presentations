"""Fill placeholders in a corporate .pptx template."""

from pathlib import Path

from loguru import logger
from pptx import Presentation

from presentations.core.schemas import DeckSpec, LayoutProfile
from presentations.llm.layout_validate import sanitize_deck_spec


def _find_placeholder(slide, ph_idx: int):
    """Find a placeholder by index on a slide."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            return shape
    return None


def compile_from_template(
    template_path: str | Path,
    deck_spec: DeckSpec,
    output_path: str | Path,
    layout_profile: LayoutProfile | None = None,
) -> Path:
    """Compile a deck by filling template placeholders.

    Args:
        template_path: Source .pptx template.
        deck_spec: Structured slide content.
        output_path: Destination .pptx path.
        layout_profile: Optional layout metadata for validation.

    Returns:
        Path to the generated presentation.
    """
    template = Path(template_path)
    destination = Path(output_path)
    destination.parent.mkdir(parents=True, exist_ok=True)

    if layout_profile is not None:
        deck_spec = sanitize_deck_spec(deck_spec, layout_profile)

    prs = Presentation(str(template))
    # Remove existing slides from template preview deck
    slide_ids = list(prs.slides._sldIdLst)  # noqa: SLF001
    for slide_id in slide_ids:
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001

    for slide_spec in deck_spec.slides:
        layout = prs.slide_layouts[slide_spec.layout_index]
        slide = prs.slides.add_slide(layout)
        for mapping in slide_spec.mappings:
            placeholder = _find_placeholder(slide, mapping.ph_idx)
            if placeholder is None:
                logger.warning(
                    "Placeholder {} not found on layout {}",
                    mapping.ph_idx,
                    slide_spec.layout_index,
                )
                continue
            placeholder.text = mapping.content

    prs.save(str(destination))
    logger.info("Compiled template deck to {}", destination)
    return destination
