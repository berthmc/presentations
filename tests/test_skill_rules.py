"""Tests for skill_rules formatting and compliance helpers."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from presentations.agents.assembler import compile_from_template
from presentations.agents.skill_rules import (
    check_leftover_placeholder_text,
    check_unicode_bullets,
    estimate_text_overflow,
    fill_placeholder_with_rules,
    normalize_content_lines,
    parse_inline_markdown,
    strip_leading_markdown,
)
from presentations.compile.theme_palette import build_theme_palette, contrast_on_white, usable_accent_colors
from presentations.core.schemas import DeckSpec, GenerationMode, PlaceholderMapping, SlideSpec


def test_normalize_content_lines_strips_unicode_bullets() -> None:
    lines = normalize_content_lines("• First item\n- Second item")
    assert lines == ["First item", "Second item"]


def test_strip_leading_markdown_removes_headers_and_numbering() -> None:
    assert strip_leading_markdown("# Heading") == "Heading"
    assert strip_leading_markdown("1. Numbered item") == "Numbered item"
    assert strip_leading_markdown("* Bullet item") == "Bullet item"


def test_parse_inline_markdown_splits_bold_segments() -> None:
    segments = parse_inline_markdown("**Transparency Rules**: Brokers must disclose.")
    texts = [segment.text for segment in segments]
    assert texts == ["Transparency Rules", ": Brokers must disclose."]
    assert segments[0].bold is True
    assert segments[0].accent is True
    assert "**" not in "".join(texts)


def test_parse_inline_markdown_splits_leadin_label_without_markers() -> None:
    segments = parse_inline_markdown("Best Execution: Firms must seek best terms.")
    assert segments[0].text == "Best Execution:"
    assert segments[0].bold is True
    assert segments[1].text == "Firms must seek best terms."


def test_fill_placeholder_with_rules_applies_theme_accent(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    placeholder = slide.placeholders[1]
    fill_placeholder_with_rules(
        placeholder,
        "**Key Point**: Supporting detail",
        accent=MSO_THEME_COLOR.ACCENT_3,
    )
    paragraph = placeholder.text_frame.paragraphs[0]
    assert len(paragraph.runs) >= 2
    assert "**" not in placeholder.text
    assert paragraph.runs[0].font.bold is True
    assert paragraph.runs[0].font.color.theme_color == MSO_THEME_COLOR.ACCENT_3


def test_check_leftover_placeholder_text_detects_lorem() -> None:
    issues = check_leftover_placeholder_text("Lorem ipsum dolor sit amet")
    assert issues


def test_check_unicode_bullets_flags_bullet_char() -> None:
    issues = check_unicode_bullets("• Item one")
    assert issues


def test_estimate_text_overflow_flags_long_content() -> None:
    text = "word " * 500
    result = estimate_text_overflow(text, width_emu=914400, height_emu=457200, font_size_pt=14.0)
    assert result["overflow"] is True
    assert result["overflow_lines"] > 0


def test_usable_accent_colors_skips_low_contrast_yellow() -> None:
    accent_hex = {
        MSO_THEME_COLOR.ACCENT_1: "44BA7E",
        MSO_THEME_COLOR.ACCENT_5: "FFD34E",
        MSO_THEME_COLOR.ACCENT_6: "DD0C86",
    }
    assert contrast_on_white("FFD34E") < 3.0
    usable = usable_accent_colors(accent_hex)
    assert MSO_THEME_COLOR.ACCENT_5 not in usable
    assert MSO_THEME_COLOR.ACCENT_6 in usable


@pytest.mark.parametrize(
    "template_path",
    [Path("documentation/briefs/i'm_an_investor.pptx")],
)
def test_template_compile_applies_accent_colours_without_markdown(
    tmp_path: Path,
    template_path: Path,
) -> None:
    if not template_path.exists():
        pytest.skip("EC template fixture not available")

    deck = DeckSpec(
        title="Palette Test",
        mode=GenerationMode.TEMPLATE,
        slides=[
            SlideSpec(
                layout_index=0,
                mappings=[
                    PlaceholderMapping(ph_idx=0, content="MiFID II Overview"),
                    PlaceholderMapping(ph_idx=10, content="Investor protection in the EU"),
                ],
            ),
            SlideSpec(
                layout_index=2,
                mappings=[
                    PlaceholderMapping(ph_idx=0, content="Transparency"),
                    PlaceholderMapping(
                        ph_idx=11,
                        content=(
                            "**Product Disclosure**: Detailed product information\n"
                            "**Risk Warning**: Clear warnings required"
                        ),
                    ),
                ],
            ),
            SlideSpec(
                layout_index=3,
                mappings=[
                    PlaceholderMapping(ph_idx=0, content="Thank you"),
                    PlaceholderMapping(ph_idx=11, content="Sources: regulatory documents"),
                ],
            ),
        ],
    )
    output = tmp_path / "palette_test.pptx"
    compile_from_template(template_path, deck, output)

    compiled = Presentation(str(output))
    palette = build_theme_palette(compiled)
    assert palette.accents

    all_text = []
    accent_found = False
    for slide in compiled.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    all_text.append(run.text)
                    try:
                        if run.font.color.theme_color is not None:
                            accent_found = True
                    except AttributeError:
                        continue

    joined = "".join(all_text)
    assert "**" not in joined
    assert accent_found
